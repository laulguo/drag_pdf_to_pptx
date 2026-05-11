#!/usr/bin/env python
"""
Convert a beamer PDF into a PowerPoint deck.

This keeps each PDF page as one full-slide image, which is the most reliable
way to preserve beamer layout, fonts, equations, and overlay steps.
"""

from __future__ import annotations

import argparse
import sys
import tempfile
from pathlib import Path


EMU_PER_INCH = 914400
POINTS_PER_INCH = 72
DEFAULT_DPI = 400


def require_dependencies() -> tuple[object, object]:
    missing: list[str] = []

    try:
        import fitz  # type: ignore
    except ImportError:
        fitz = None
        missing.append("PyMuPDF")

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        Presentation = None
        missing.append("python-pptx")

    if missing:
        names = " ".join(missing)
        print("Missing dependencies:", ", ".join(missing), file=sys.stderr)
        print(f"Please run: python -m pip install {names}", file=sys.stderr)
        raise SystemExit(2)

    return fitz, Presentation


def slide_size_from_pdf_page(page: object) -> tuple[int, int]:
    rect = page.rect
    width_in = rect.width / POINTS_PER_INCH
    height_in = rect.height / POINTS_PER_INCH
    return int(width_in * EMU_PER_INCH), int(height_in * EMU_PER_INCH)


def convert_pdf_to_pptx(
    pdf_path: Path,
    output_path: Path | None = None,
    dpi: int = DEFAULT_DPI,
    overwrite: bool = True,
) -> Path:
    fitz, Presentation = require_dependencies()

    pdf_path = pdf_path.expanduser().resolve()
    if not pdf_path.exists():
        raise FileNotFoundError(f"File not found: {pdf_path}")
    if pdf_path.suffix.lower() != ".pdf":
        raise ValueError(f"Only PDF files are supported: {pdf_path}")

    if output_path is None:
        output_path = pdf_path.with_suffix(".pptx")
    else:
        output_path = output_path.expanduser().resolve()

    if output_path.exists() and not overwrite:
        raise FileExistsError(f"Output file already exists: {output_path}")

    doc = fitz.open(str(pdf_path))
    if doc.page_count == 0:
        raise ValueError(f"PDF has no pages: {pdf_path}")

    prs = Presentation()
    prs.slide_width, prs.slide_height = slide_size_from_pdf_page(doc[0])
    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory(prefix="pdf_to_pptx_") as temp_dir:
        temp_path = Path(temp_dir)

        for index, page in enumerate(doc, start=1):
            image_path = temp_path / f"page_{index:04d}.png"
            pixmap = page.get_pixmap(dpi=dpi, alpha=False)
            pixmap.save(str(image_path))

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            print(f"Converted {index}/{doc.page_count} pages")

    prs.save(str(output_path))
    doc.close()
    return output_path


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Convert a beamer PDF to PPTX. Each PDF page becomes one PowerPoint slide."
    )
    parser.add_argument("pdf", nargs="+", help="PDF file path(s). Multiple files are supported.")
    parser.add_argument(
        "-o",
        "--output",
        help="Output PPTX path. Only available for one PDF. Defaults to same folder and name as the PDF.",
    )
    parser.add_argument(
        "--dpi",
        type=int,
        default=DEFAULT_DPI,
        help=f"Render DPI. Default: {DEFAULT_DPI}. Use 200 for smaller files, or 600 for sharper slides.",
    )
    parser.add_argument(
        "--no-overwrite",
        action="store_true",
        help="Stop if the output file already exists instead of overwriting it.",
    )
    return parser.parse_args(argv)


def main(argv: list[str] | None = None) -> int:
    args = parse_args(sys.argv[1:] if argv is None else argv)

    pdf_paths = [Path(item) for item in args.pdf]
    if args.output and len(pdf_paths) > 1:
        print("--output cannot be used when converting multiple PDFs.", file=sys.stderr)
        return 2

    try:
        for pdf_path in pdf_paths:
            output_path = Path(args.output) if args.output else None
            result = convert_pdf_to_pptx(
                pdf_path,
                output_path=output_path,
                dpi=args.dpi,
                overwrite=not args.no_overwrite,
            )
            print(f"Done: {result}")
    except Exception as exc:
        print(f"Conversion failed: {exc}", file=sys.stderr)
        return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
