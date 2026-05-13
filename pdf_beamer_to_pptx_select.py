#!/usr/bin/env python
"""
Convert a PDF in the current directory into a PowerPoint deck.

This version is designed for IDE usage such as Spyder:
- No drag-and-drop or command-line arguments required
- Lists PDF files in the current directory
- Lets you choose one by number
"""

from __future__ import annotations

import sys
import tempfile
from pathlib import Path


EMU_PER_INCH = 914400
POINTS_PER_INCH = 72
DEFAULT_DPI = 400


def require_dependencies() -> tuple[object, object]:
    missing: list[str] = []

    try:
        import fitz  # type: ignore
    except ImportError:
        fitz = None
        missing.append("PyMuPDF")

    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        Presentation = None
        missing.append("python-pptx")

    if missing:
        names = " ".join(missing)
        print("Missing dependencies:", ", ".join(missing), file=sys.stderr)
        print(f"Please run: python -m pip install {names}", file=sys.stderr)
        raise SystemExit(2)

    return fitz, Presentation


def slide_size_from_pdf_page(page: object) -> tuple[int, int]:
    rect = page.rect
    width_in = rect.width / POINTS_PER_INCH
    height_in = rect.height / POINTS_PER_INCH
    return int(width_in * EMU_PER_INCH), int(height_in * EMU_PER_INCH)


def convert_pdf_to_pptx(
    pdf_path: Path,
    output_path: Path | None = None,
    dpi: int = DEFAULT_DPI,
    overwrite: bool = True,
) -> Path:
    fitz, Presentation = require_dependencies()

    pdf_path = pdf_path.expanduser().resolve()
    if not pdf_path.exists():
        raise FileNotFoundError(f"File not found: {pdf_path}")
    if pdf_path.suffix.lower() != ".pdf":
        raise ValueError(f"Only PDF files are supported: {pdf_path}")

    if output_path is None:
        output_path = pdf_path.with_suffix(".pptx")
    else:
        output_path = output_path.expanduser().resolve()

    if output_path.exists() and not overwrite:
        raise FileExistsError(f"Output file already exists: {output_path}")

    doc = fitz.open(str(pdf_path))
    if doc.page_count == 0:
        raise ValueError(f"PDF has no pages: {pdf_path}")

    prs = Presentation()
    prs.slide_width, prs.slide_height = slide_size_from_pdf_page(doc[0])
    blank_layout = prs.slide_layouts[6]

    with tempfile.TemporaryDirectory(prefix="pdf_to_pptx_") as temp_dir:
        temp_path = Path(temp_dir)

        for index, page in enumerate(doc, start=1):
            image_path = temp_path / f"page_{index:04d}.png"
            pixmap = page.get_pixmap(dpi=dpi, alpha=False)
            pixmap.save(str(image_path))

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            print(f"Converted {index}/{doc.page_count} pages")

    prs.save(str(output_path))
    doc.close()
    return output_path


def find_pdf_files(base_dir: Path) -> list[Path]:
    return sorted(
        [path for path in base_dir.iterdir() if path.is_file() and path.suffix.lower() == ".pdf"],
        key=lambda path: path.name.lower(),
    )


def choose_pdf(pdf_files: list[Path]) -> Path:
    while True:
        print("\nAvailable PDF files:")
        for index, pdf_path in enumerate(pdf_files, start=1):
            print(f"  {index}. {pdf_path.name}")

        choice = input("\nEnter the PDF number to convert: ").strip()
        if not choice:
            print("Please enter a number.")
            continue
        if not choice.isdigit():
            print("Please enter digits only.")
            continue

        index = int(choice)
        if 1 <= index <= len(pdf_files):
            return pdf_files[index - 1]

        print(f"Please enter a number between 1 and {len(pdf_files)}.")


def ask_overwrite(output_path: Path) -> bool:
    if not output_path.exists():
        return True

    while True:
        choice = input(f"\n{output_path.name} already exists. Overwrite it? (y/n): ").strip().lower()
        if choice in {"y", "yes"}:
            return True
        if choice in {"n", "no"}:
            return False
        print("Please enter y or n.")


def main() -> int:
    base_dir = Path.cwd()
    print(f"Current directory: {base_dir}")

    pdf_files = find_pdf_files(base_dir)
    if not pdf_files:
        print("No PDF files were found in the current directory.", file=sys.stderr)
        return 1

    pdf_path = choose_pdf(pdf_files)
    output_path = pdf_path.with_suffix(".pptx")

    print(f"\nSelected PDF: {pdf_path.name}")
    print(f"Output PPTX: {output_path.name}")

    overwrite = ask_overwrite(output_path)
    if not overwrite:
        print("Cancelled. No file was written.")
        return 0

    try:
        result = convert_pdf_to_pptx(
            pdf_path,
            output_path=output_path,
            dpi=DEFAULT_DPI,
            overwrite=True,
        )
    except Exception as exc:
        print(f"Conversion failed: {exc}", file=sys.stderr)
        return 1

    print(f"\nDone: {result}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
